"""Office-document modality readers — email, presentations, spreadsheets.

CONCEPT:KG-2.66 — multi-modal ingestion readers (office tail).

Extends the universal ingestion funnel's *reader* stage (reader → structure-router
→ {open|schema} extraction → ontology grounding) to the everyday office formats
the PDF/DOCX/EPUB readers don't cover:

* **email**   ``.eml`` (stdlib :mod:`email`) and ``.msg`` (optional ``extract-msg``)
  → ``From/To/Cc/Subject/Date`` headers followed by the plain-text body.
* **slides**  ``.pptx`` (optional ``python-pptx``) → per-slide text (shapes, tables,
  speaker notes), one block per slide.
* **sheets**  ``.xlsx`` (optional ``openpyxl``) and ``.csv``/``.tsv`` (stdlib
  :mod:`csv`) → ``# Sheet`` headers followed by tab-joined rows.

Each reader is a ``(path: str) -> str`` function self-registered against the
sibling reader registry (``extraction.readers.register_reader``) keyed by its
extensions, so the structure-router picks them up with no edits to a shared hub
(the same self-registration pattern as ``protocols/source_connectors`` KG-2.59
and ``enrichment/registry`` KG-2.9).

Discipline (matches ``kb/parser`` / ``enrichment/extractors/document``):

* every heavy dependency is **import-guarded and auto-detected** — a reader runs
  its full extraction when the lib is importable and degrades to a clear,
  logged no-op (returns ``""``) when it is not, so ingestion never raises;
* the readers are **best-effort** in the ingest path — any per-file failure is
  caught and logged, never propagated.
"""

from __future__ import annotations

import csv
import logging
from collections.abc import Callable
from pathlib import Path

logger = logging.getLogger(__name__)

# Reader signature: an on-disk path → extracted plain text.
Reader = Callable[[str], str]

# ----------------------------------------------------------------------------
# Reader-registry binding (sibling component).
#
# ``extraction.readers`` owns ``register_reader`` — a decorator taking one or
# more extensions and registering the wrapped ``(path) -> str`` reader. We import
# it best-effort: if the registry module is not present yet, fall back to a
# no-op decorator so this module still imports cleanly and its readers remain
# directly callable (the structure-router simply won't auto-discover them until
# the registry lands). This mirrors the guarded-optional-import convention used
# throughout the ingestion package.
# ----------------------------------------------------------------------------
try:
    from .readers import register_reader  # type: ignore[attr-defined]
except Exception:  # pragma: no cover - registry not present yet

    def register_reader(*extensions: str) -> Callable[[Reader], Reader]:
        """Fallback no-op decorator used until ``extraction.readers`` exists."""

        def _decorator(fn: Reader) -> Reader:
            return fn

        return _decorator


def _read_bytes(path: str) -> bytes:
    return Path(path).read_bytes()


# ============================================================================
# Email — .eml (stdlib) / .msg (optional extract-msg)
# ============================================================================

_EMAIL_HEADER_FIELDS = ("From", "To", "Cc", "Subject", "Date")


def _format_email(headers: dict[str, str], body: str, *, source: str = "") -> str:
    """Render a header block (only the populated fields) followed by the body."""
    lines = [
        f"{field}: {headers[field]}"
        for field in _EMAIL_HEADER_FIELDS
        if headers.get(field)
    ]
    head = "\n".join(lines)
    body = (body or "").strip()
    if head and body:
        return f"{head}\n\n{body}"
    return head or body


@register_reader(".eml")
def read_eml(path: str) -> str:
    """Extract an ``.eml`` message via the stdlib :mod:`email` parser.

    Returns ``From/To/Cc/Subject/Date`` headers followed by the plain-text body
    (HTML-only mails are de-tagged best-effort). Stdlib only — always available.
    """
    from email import policy
    from email.parser import BytesParser

    try:
        msg = BytesParser(policy=policy.default).parsebytes(_read_bytes(path))
    except Exception as exc:  # pragma: no cover - unreadable file
        logger.debug("read_eml failed for %s: %s", path, exc)
        return ""

    headers = {f: str(msg.get(f, "") or "").strip() for f in _EMAIL_HEADER_FIELDS}
    body = _extract_email_body(msg)
    return _format_email(headers, body, source=path)


def _extract_email_body(msg: object) -> str:
    """Pull the best text body out of a parsed :class:`email.message.Message`."""
    get_body = getattr(msg, "get_body", None)
    if callable(get_body):  # modern EmailMessage API
        try:
            part = get_body(preferencelist=("plain", "html"))
            if part is not None:
                content = part.get_content()
                if part.get_content_subtype() == "html":
                    return _strip_html(content)
                return content
        except Exception:  # pragma: no cover - exotic payloads
            pass
    # Fallback: walk parts.
    texts: list[str] = []
    for part in getattr(msg, "walk", lambda: [])():
        ctype = part.get_content_type()
        if ctype == "text/plain":
            with _suppress():
                texts.append(part.get_content())
        elif ctype == "text/html" and not texts:
            with _suppress():
                texts.append(_strip_html(part.get_content()))
    return "\n".join(t for t in texts if t)


@register_reader(".msg")
def read_msg(path: str) -> str:
    """Extract an Outlook ``.msg`` via the optional ``extract-msg`` library.

    Auto-detected: runs when ``extract_msg`` is importable, otherwise degrades to
    a clear no-op (returns ``""``) so ingestion never raises.
    """
    try:
        import extract_msg  # type: ignore[import-not-found]
    except ImportError:
        logger.warning(
            "extract-msg not installed; cannot read %s "
            "(install with: pip install extract-msg). Skipping.",
            path,
        )
        return ""
    try:
        m = extract_msg.Message(path)
        headers = {
            "From": str(getattr(m, "sender", "") or ""),
            "To": str(getattr(m, "to", "") or ""),
            "Cc": str(getattr(m, "cc", "") or ""),
            "Subject": str(getattr(m, "subject", "") or ""),
            "Date": str(getattr(m, "date", "") or ""),
        }
        body = str(getattr(m, "body", "") or "")
        try:
            m.close()
        except Exception:  # pragma: no cover
            pass
        return _format_email(headers, body, source=path)
    except Exception as exc:  # pragma: no cover - corrupt .msg
        logger.debug("read_msg failed for %s: %s", path, exc)
        return ""


# ============================================================================
# Presentations — .pptx (optional python-pptx)
# ============================================================================


@register_reader(".pptx")
def read_pptx(path: str) -> str:
    """Extract slide text from a ``.pptx`` via the optional ``python-pptx`` library.

    One text block per slide: shape text, table cells, and speaker notes. Slides
    are separated by a ``--- Slide N ---`` marker so downstream chunking keeps
    slide boundaries. Auto-detected; degrades to a clear no-op if ``pptx`` is
    absent.
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        logger.warning(
            "python-pptx not installed; cannot read %s "
            "(install with: pip install python-pptx). Skipping.",
            path,
        )
        return ""
    try:
        prs = Presentation(path)
    except Exception as exc:  # pragma: no cover - corrupt pptx
        logger.debug("read_pptx failed to open %s: %s", path, exc)
        return ""

    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            _collect_shape_text(shape, parts)
        notes = _pptx_notes(slide)
        if notes:
            parts.append(f"[notes] {notes}")
        body = "\n".join(p for p in parts if p.strip())
        if body.strip():
            slides.append(f"--- Slide {idx} ---\n{body}")
    return "\n\n".join(slides)


def _collect_shape_text(shape: object, parts: list[str]) -> None:
    """Append text from a pptx shape (text frame or table) into ``parts``."""
    if getattr(shape, "has_text_frame", False):
        text = (getattr(shape, "text", "") or "").strip()
        if text:
            parts.append(text)
    if getattr(shape, "has_table", False):
        try:
            table = shape.table
            for row in table.rows:
                cells = [(c.text or "").strip() for c in row.cells]
                line = "\t".join(c for c in cells)
                if line.strip():
                    parts.append(line)
        except Exception:  # pragma: no cover - exotic table
            pass


def _pptx_notes(slide: object) -> str:
    """Speaker notes for a slide, or empty string."""
    try:
        if getattr(slide, "has_notes_slide", False):
            tf = slide.notes_slide.notes_text_frame
            return (getattr(tf, "text", "") or "").strip()
    except Exception:  # pragma: no cover
        pass
    return ""


# ============================================================================
# Spreadsheets — .xlsx (optional openpyxl) / .csv / .tsv (stdlib)
# ============================================================================


def _format_rows(rows: list[list[str]]) -> str:
    """Tab-join each non-empty row into a line block."""
    lines = []
    for row in rows:
        cells = ["" if c is None else str(c) for c in row]
        line = "\t".join(cells).rstrip("\t")
        if line.strip():
            lines.append(line)
    return "\n".join(lines)


@register_reader(".xlsx")
def read_xlsx(path: str) -> str:
    """Extract sheet/row text from an ``.xlsx`` via the optional ``openpyxl`` library.

    One ``# <sheet>`` header per worksheet followed by tab-joined rows (values
    only, ``read_only`` + ``data_only`` for speed and to surface computed cell
    values). Auto-detected; degrades to a clear no-op if ``openpyxl`` is absent.
    """
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except ImportError:
        logger.warning(
            "openpyxl not installed; cannot read %s "
            "(install with: pip install openpyxl). Skipping.",
            path,
        )
        return ""
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:  # pragma: no cover - corrupt xlsx
        logger.debug("read_xlsx failed to open %s: %s", path, exc)
        return ""

    blocks: list[str] = []
    try:
        for ws in wb.worksheets:
            rows = [list(r) for r in ws.iter_rows(values_only=True)]
            body = _format_rows(rows)
            if body.strip():
                blocks.append(f"# {ws.title}\n{body}")
    finally:
        with _suppress():
            wb.close()
    return "\n\n".join(blocks)


def read_csv(path: str) -> str:
    """Extract row text from a ``.csv``/``.tsv`` via the stdlib :mod:`csv` reader.

    Not registered for ``.csv``/``.tsv`` — the sibling ``extraction.readers``
    module owns those extensions (one source of truth); this richer
    delimiter-sniffing variant remains importable for direct use.

    Delimiter is sniffed (falls back to ``,`` for ``.csv`` and ``\\t`` for
    ``.tsv``); rows are re-emitted tab-joined. Stdlib only — always available.
    """
    try:
        raw = Path(path).read_text(encoding="utf-8", errors="replace")
    except OSError as exc:  # pragma: no cover - unreadable file
        logger.debug("read_csv failed for %s: %s", path, exc)
        return ""
    if not raw.strip():
        return ""

    default = "\t" if path.lower().endswith(".tsv") else ","
    sample = raw[:8192]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
        delimiter = dialect.delimiter
    except csv.Error:
        delimiter = default

    rows: list[list[str]] = []
    try:
        reader = csv.reader(raw.splitlines(), delimiter=delimiter)
        rows = [list(r) for r in reader]
    except csv.Error as exc:  # pragma: no cover - malformed csv
        logger.debug("read_csv parse failed for %s: %s", path, exc)
        return raw
    return _format_rows(rows)


# ============================================================================
# Small shared helpers
# ============================================================================


def _strip_html(html: str) -> str:
    """De-tag HTML best-effort (BeautifulSoup if present, else regex)."""
    try:
        from bs4 import BeautifulSoup

        return BeautifulSoup(html, "html.parser").get_text(separator="\n", strip=True)
    except ImportError:
        import re

        return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", html)).strip()


class _suppress:
    """Tiny ``contextlib.suppress(Exception)`` stand-in (avoids the import churn)."""

    def __enter__(self) -> None:
        return None

    def __exit__(self, exc_type: object, exc: object, tb: object) -> bool:
        return True
