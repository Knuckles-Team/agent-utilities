from __future__ import annotations

"""Self-registering multi-modal READER REGISTRY (CONCEPT:KG-2.66 family).

The universal-ingestion funnel needs ONE text-extraction front door so every new
modality (slides / spreadsheets / audio / scanned images / html …) flows through
the *same* downstream pipeline (structure-router → open|schema extraction →
ontology grounding → OWL-RL/SHACL closure) as plain documents do.

This module is that front door. It mirrors the project's other self-registering
decorator + ``pkgutil`` discovery pattern
(:mod:`agent_utilities.protocols.source_connectors.registry`):

* ``@register_reader(*exts)`` — class/function decorator that maps one or more
  lowercase file extensions (``".pptx"``) and/or MIME types
  (``"application/pdf"``) to a reader callable ``(path: str) -> str``.
* ``_READER_REGISTRY`` — the extension/MIME → reader map.
* ``discover()`` — imports the ``readers`` subpackage so every decorator runs,
  exactly like the connector registry imports ``connectors``. Cached + idempotent.
* ``read_any(path)`` — the dispatcher. Resolves a reader by extension (then MIME),
  and **falls back to the existing**
  :func:`agent_utilities.knowledge_graph.enrichment.extractors.document.read_document_text`
  for ``.pdf`` / ``.md`` / ``.txt`` / ``.rst`` / ``.json`` / ``.eml`` so the
  battle-tested (PyMuPDF-fast) path stays the source of truth for those.

A handful of built-in readers are registered **inline** below (csv/tsv, html,
pptx, xlsx, audio, image-OCR). ``.txt`` / ``.md`` / ``.pdf`` need no registration
— they are served by the fallback. Heavy deps (python-pptx, openpyxl,
faster-whisper, pytesseract) are **import-guarded and auto-detected**: a reader
runs when its dependency imports and degrades to a clear no-op (empty string,
logged once) otherwise — honouring the ingest contract of *best-effort, never
raise*.

Wire-First note (see AGENTS.md): ``check_wiring.py`` cannot see decorator
registration, so :func:`discover` is invoked on the **live** ``read_any`` path
(every ``read_any`` call calls ``discover`` first) and asserted by a live-path
test — it is not relied on for import-graph reachability. The ingestion engine /
document adaptor calls :func:`read_any`, which routes new modalities through the
same pipeline (see the wiring note).
"""

import importlib
import logging
import os
import pkgutil
from collections.abc import Callable
from typing import TypeVar

logger = logging.getLogger(__name__)

__all__ = [
    "Reader",
    "register_reader",
    "discover",
    "get_reader",
    "list_readers",
    "read_any",
]

#: A reader turns a path into its extracted text (best-effort, never raises).
Reader = Callable[[str], str]

_READER_REGISTRY: dict[str, Reader] = {}
_DISCOVERED = False

R = TypeVar("R", bound=Reader)

# Extensions whose canonical reader is the existing ``read_document_text`` fast
# path (PyMuPDF for PDF, direct decode for the text family). We do NOT register
# these — ``read_any`` routes them to the fallback so there is one source of
# truth for the document path.
_FALLBACK_EXTS = frozenset({".pdf", ".md", ".txt", ".rst", ".json", ".eml"})

# Minimal extension → MIME map so a caller that only knows a MIME type (e.g. an
# HTTP download with no filename) can still dispatch. Kept tiny on purpose —
# ``mimetypes`` covers the long tail at call time.
_MIME_TO_EXT = {
    "text/csv": ".csv",
    "text/tab-separated-values": ".tsv",
    "text/html": ".html",
    "application/xhtml+xml": ".html",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",  # noqa: E501
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.ms-excel": ".xls",
    "audio/mpeg": ".mp3",
    "audio/wav": ".wav",
    "audio/x-wav": ".wav",
    "audio/mp4": ".m4a",
    "audio/flac": ".flac",
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/tiff": ".tiff",
}


def register_reader(*exts: str) -> Callable[[R], R]:
    """Decorator registering a reader under one or more extensions/MIME types.

    CONCEPT:KG-2.66. Keys are normalised lowercase; an extension is stored with a
    leading dot (``"csv"`` and ``".csv"`` both register under ``".csv"``); a MIME
    type (contains ``"/"``) is mapped through :data:`_MIME_TO_EXT` to its
    extension so all dispatch is extension-keyed. Idempotent — re-registering the
    same key overwrites, so repeated :func:`discover` calls are safe.

    Example::

        @register_reader(".csv", ".tsv")
        def read_delimited(path: str) -> str:
            ...
    """

    def _decorator(fn: R) -> R:
        for raw in exts:
            key = _normalize_key(raw)
            if key is None:
                continue
            _READER_REGISTRY[key] = fn
            logger.debug(
                "[KG-2.66] registered reader %r -> %s", key, getattr(fn, "__name__", fn)
            )
        return fn

    return _decorator


def _normalize_key(raw: str) -> str | None:
    """Normalise an extension or MIME type to a dotted lowercase extension key."""
    if not raw:
        return None
    raw = raw.strip().lower()
    if "/" in raw:  # a MIME type
        return _MIME_TO_EXT.get(raw)
    return raw if raw.startswith(".") else f".{raw}"


#: Sibling reader modules whose import triggers their ``@register_reader``
#: decorators. The built-in readers ship as ``readers_*`` modules next to this
#: registry (not a ``readers/`` subpackage), so discovery imports them directly;
#: an optional ``readers/`` plugin subpackage is also swept if present.
_BUILTIN_READER_MODULES = ("readers_office", "readers_media")


def discover() -> dict[str, Reader]:
    """Import the built-in reader modules (and any ``readers`` plugin subpackage)
    so every ``@register_reader`` decorator runs; return the registry map.

    CONCEPT:KG-2.66 — the discovery half of the self-registering pattern, called
    on the live :func:`read_any` path. Idempotent: module imports are cached and
    registration overwrites, so repeated calls are safe.
    """
    global _DISCOVERED
    if _DISCOVERED:
        return dict(_READER_REGISTRY)
    for name in _BUILTIN_READER_MODULES:
        try:
            importlib.import_module(f"{__package__}.{name}")
        except Exception as exc:  # noqa: BLE001 — one bad reader must not block others
            logger.warning("[KG-2.66] failed to import reader module %r: %s", name, exc)
    # Optional ``readers/`` plugin subpackage (a directory) for out-of-tree readers.
    try:
        pkg = importlib.import_module(f"{__package__}.readers")
        if hasattr(pkg, "__path__"):
            for mod in pkgutil.iter_modules(pkg.__path__):
                if mod.name.startswith("_"):
                    continue
                try:
                    importlib.import_module(f"{pkg.__name__}.{mod.name}")
                except Exception as exc:  # noqa: BLE001
                    logger.warning(
                        "[KG-2.66] failed to import reader %r: %s", mod.name, exc
                    )
    except ModuleNotFoundError:
        pass
    except Exception as exc:  # noqa: BLE001 — discovery must not break ingestion
        logger.warning("[KG-2.66] reader discovery failed: %s", exc)
    _DISCOVERED = True
    return dict(_READER_REGISTRY)


def get_reader(key: str) -> Reader | None:
    """Return the reader registered for an extension or MIME type (after discovery)."""
    discover()
    norm = _normalize_key(key)
    return _READER_REGISTRY.get(norm) if norm else None


def list_readers() -> list[str]:
    """Return every registered extension key (after discovery), sorted."""
    discover()
    return sorted(_READER_REGISTRY)


def read_any(path: str, *, mime: str | None = None) -> str:
    """Extract text from ``path`` through the registry, with the document fallback.

    CONCEPT:KG-2.66 — the universal text front door. Dispatch order:

    1. The existing :func:`read_document_text` fast path for the document family
       (``.pdf`` / ``.md`` / ``.txt`` / ``.rst`` / ``.json`` / ``.eml``) — one
       source of truth, PyMuPDF-fast.
    2. A registered reader matched by the file extension (then by ``mime`` /
       :pyfunc:`mimetypes.guess_type` when the path has no useful suffix).
    3. The :func:`read_document_text` fallback for anything else (it returns
       "" for formats it can't read, which is the correct best-effort no-op).

    Best-effort and **never raises** — an unreadable/absent file or a missing
    optional dependency yields ``""`` (skipped upstream), matching the ingest
    contract.
    """
    discover()
    ext = os.path.splitext(path)[1].lower()

    # 1) Document family → existing fast path (source of truth).
    if ext in _FALLBACK_EXTS:
        return _fallback_read(path)

    # 2) Registered modality reader (by extension, then by MIME).
    reader = _READER_REGISTRY.get(ext)
    if reader is None:
        guessed = _normalize_key(mime) if mime else None
        if guessed is None and not ext:
            import mimetypes

            guess, _ = mimetypes.guess_type(path)
            guessed = _normalize_key(guess) if guess else None
        if guessed is not None:
            reader = _READER_REGISTRY.get(guessed)
    if reader is not None:
        try:
            return reader(path) or ""
        except Exception as exc:  # noqa: BLE001 — reader failure → empty (skipped)
            logger.warning("[KG-2.66] reader for %r failed on %s: %s", ext, path, exc)
            return ""

    # 3) Unknown modality → let the document reader try (returns "" if it can't).
    return _fallback_read(path)


def _fallback_read(path: str) -> str:
    """Delegate to the canonical ``read_document_text`` (lazy import, never raises)."""
    try:
        from ..enrichment.extractors.document import read_document_text

        return read_document_text(path)
    except Exception as exc:  # noqa: BLE001 — best-effort
        logger.debug("[KG-2.66] fallback read failed for %s: %s", path, exc)
        return ""


# --------------------------------------------------------------------------- #
# Built-in inline readers. Each is best-effort and import-guarded: it runs when
# its (optional) dependency is importable and degrades to a clear no-op when not.
# These cover the common always-available formats (csv/tsv, html via stdlib)
# plus the heavy auto-detected ones (pptx, xlsx, audio, image-OCR).
# --------------------------------------------------------------------------- #


@register_reader(".csv", ".tsv")
def read_delimited(path: str) -> str:
    """CSV/TSV → newline-joined rows, cells separated by " | " (stdlib only).

    A flat textual rendering is intentional: it feeds the same concept/fact
    extraction the document path uses, no schema inference here.
    """
    import csv

    delim = "\t" if path.lower().endswith(".tsv") else ","
    try:
        with open(path, newline="", encoding="utf-8", errors="ignore") as fh:
            rows = [
                " | ".join(cell.strip() for cell in row)
                for row in csv.reader(fh, delimiter=delim)
            ]
        return "\n".join(r for r in rows if r.strip())
    except OSError:
        return ""


@register_reader(".html", ".htm", ".xhtml")
def read_html(path: str) -> str:
    """HTML → visible text. Prefers BeautifulSoup when present, else a stdlib
    ``HTMLParser`` that strips script/style and collapses whitespace."""
    try:
        raw = open(path, encoding="utf-8", errors="ignore").read()
    except OSError:
        return ""
    try:
        from bs4 import BeautifulSoup  # type: ignore[import-not-found]

        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "head", "meta", "noscript"]):
            tag.decompose()
        return "\n".join(
            line for line in soup.get_text("\n").splitlines() if line.strip()
        )
    except ImportError:
        pass
    from html.parser import HTMLParser

    class _Text(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self._skip = 0
            self.parts: list[str] = []

        def handle_starttag(self, tag: str, attrs: object) -> None:
            if tag in ("script", "style", "head"):
                self._skip += 1

        def handle_endtag(self, tag: str) -> None:
            if tag in ("script", "style", "head") and self._skip:
                self._skip -= 1

        def handle_data(self, data: str) -> None:
            if not self._skip and data.strip():
                self.parts.append(data.strip())

    parser = _Text()
    try:
        parser.feed(raw)
    except Exception:  # noqa: BLE001 — malformed HTML → whatever we parsed
        pass
    return "\n".join(parser.parts)


@register_reader(".pptx")
def read_pptx(path: str) -> str:
    """PowerPoint → slide text (auto-detected: needs ``python-pptx``).

    No-op (empty string, logged once) when ``python-pptx`` is not installed.
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError:
        logger.info("[KG-2.66] .pptx reader skipped (install python-pptx to enable)")
        return ""
    try:
        prs = Presentation(path)
    except Exception as exc:  # noqa: BLE001
        logger.warning("[KG-2.66] python-pptx failed on %s: %s", path, exc)
        return ""
    out: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        out.append(f"# Slide {idx}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                out.append(text.strip())
        notes = getattr(slide, "notes_slide", None)
        if notes is not None:
            note_text = getattr(getattr(notes, "notes_text_frame", None), "text", "")
            if note_text and note_text.strip():
                out.append(f"[notes] {note_text.strip()}")
    return "\n".join(out)


@register_reader(".xlsx", ".xlsm")
def read_xlsx(path: str) -> str:
    """Excel workbook → per-sheet textual rows (auto-detected: needs ``openpyxl``).

    Reads values only (``data_only=True``) and degrades to a no-op when
    ``openpyxl`` is absent.
    """
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except ImportError:
        logger.info("[KG-2.66] .xlsx reader skipped (install openpyxl to enable)")
        return ""
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        logger.warning("[KG-2.66] openpyxl failed on %s: %s", path, exc)
        return ""
    out: list[str] = []
    try:
        for ws in wb.worksheets:
            out.append(f"# Sheet {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [
                    str(c).strip() for c in row if c is not None and str(c).strip()
                ]
                if cells:
                    out.append(" | ".join(cells))
    finally:
        wb.close()
    return "\n".join(out)


@register_reader(".mp3", ".wav", ".m4a", ".flac", ".ogg", ".mp4", ".webm")
def read_audio(path: str) -> str:
    """Audio/video → transcript (auto-detected: needs ``faster-whisper``).

    Transcribes with a small CPU-friendly model so the resulting text flows
    through the same concept/fact pipeline. No-op when ``faster-whisper`` is not
    installed.
    """
    try:
        from faster_whisper import WhisperModel  # type: ignore[import-not-found]
    except ImportError:
        logger.info("[KG-2.66] audio reader skipped (install faster-whisper to enable)")
        return ""
    try:
        model = WhisperModel("base", device="cpu", compute_type="int8")
        segments, _info = model.transcribe(path)
        return " ".join(seg.text.strip() for seg in segments if seg.text.strip())
    except Exception as exc:  # noqa: BLE001
        logger.warning("[KG-2.66] faster-whisper failed on %s: %s", path, exc)
        return ""


@register_reader(".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif")
def read_image_ocr(path: str) -> str:
    """Scanned image → OCR text (auto-detected: needs ``pytesseract`` + Pillow).

    No-op when either ``pytesseract``/``Pillow`` is missing or the tesseract
    binary is unavailable.
    """
    try:
        import pytesseract  # type: ignore[import-not-found]
        from PIL import Image  # type: ignore[import-not-found]
    except ImportError:
        logger.info(
            "[KG-2.66] image-OCR reader skipped (install pytesseract+Pillow to enable)"
        )
        return ""
    try:
        with Image.open(path) as img:
            return (pytesseract.image_to_string(img) or "").strip()
    except Exception as exc:  # noqa: BLE001 — missing tesseract binary, bad image, …
        logger.warning("[KG-2.66] pytesseract failed on %s: %s", path, exc)
        return ""
