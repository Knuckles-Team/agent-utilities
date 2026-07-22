"""Live-path proof for the ``PageBox`` evidence-locus producer
(CONCEPT:AU-KG.identity.evidence-spine-convergence, Evidence seam completion).

Pass 1's re-check found ``store_document_page_evidence`` had ZERO external
callers despite being the very first, template locus of this seam — only its
own unit test called it (``docs/architecture/evidence_spine_convergence.md``).
This pass surveyed AU's document-page extraction surface: the PDF text path
(``extraction/pdf.py``) deliberately isolates parsing in a killable,
engine-less subprocess and joins ALL pages into one flat string (no page
boundary survives the pipe by design, for security isolation) — so it has no
computed-and-discarded per-page box to wire without adding new, unbounded
computation. ``readers_office.py``'s ``.pptx`` reader (``read_pptx``) DOES
already loop per slide (page-equivalent) with real, already-known slide
dimensions (``Presentation.slide_width``/``slide_height``) — this test proves
that gap is now closed: the EXISTING pptx reader entry point (reached from
every ``.pptx`` ingested through ``read_media``/``read_any``) now also writes a
``PageBox`` evidence locus per slide, as a side effect.

Two layers, matching ``test_table_cell_evidence_live_path``'s convention:
:func:`test_persist_slide_page_evidence_writes_the_full_slide_extent` exercises
the wiring function directly (dependency-free, runs everywhere);
:func:`test_read_pptx_writes_document_page_evidence` goes through the REAL
``read_pptx`` + a real ``python-pptx`` deck end-to-end (skipped when
``python-pptx`` isn't installed).
"""

from __future__ import annotations

import importlib
import sys

import pytest

from agent_utilities.knowledge_graph.extraction import readers_office as ro
from agent_utilities.knowledge_graph.memory import native_ingest


class _FakeStore:
    def __init__(self) -> None:
        self.page_calls: list[tuple[bytes, dict]] = []

    def store_document_page_evidence(self, data: bytes, **kwargs):
        self.page_calls.append((data, kwargs))
        return object()


def test_persist_slide_page_evidence_writes_the_full_slide_extent(monkeypatch):
    store = _FakeStore()
    monkeypatch.setattr(native_ingest, "media_store", lambda: store)

    ro._persist_slide_page_evidence(
        "/data/deck.pptx", 3, "Quarterly Results\nRevenue up 12%", 9144000, 6858000
    )

    assert len(store.page_calls) == 1
    data, kw = store.page_calls[0]
    assert data == b"Quarterly Results\nRevenue up 12%"
    assert kw["document_id"] == "/data/deck.pptx"
    assert kw["page"] == 3
    assert kw["x"] == 0.0 and kw["y"] == 0.0
    assert kw["width"] == 9144000.0
    assert kw["height"] == 6858000.0
    assert kw["source"] == "pptx"


def test_persist_slide_page_evidence_skips_empty_slide(monkeypatch):
    store = _FakeStore()
    monkeypatch.setattr(native_ingest, "media_store", lambda: store)
    ro._persist_slide_page_evidence("/data/deck.pptx", 1, "   ", 9144000, 6858000)
    assert store.page_calls == []


@pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None,
    reason="python-pptx not installed",
)
def test_read_pptx_writes_document_page_evidence(
    tmp_path, monkeypatch
):  # pragma: no cover - dep-gated
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello Slide"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))

    store = _FakeStore()
    monkeypatch.setattr(native_ingest, "media_store", lambda: store)

    # The EXISTING pptx reader entry point — no new call site.
    text = ro.read_pptx(str(p))

    # Text extraction is unchanged.
    assert "Hello Slide" in text

    # AND a PageBox evidence locus was written for the slide's full extent.
    assert len(store.page_calls) == 1
    _data, kw = store.page_calls[0]
    assert kw["document_id"] == str(p)
    assert kw["page"] == 1
    assert kw["x"] == 0.0 and kw["y"] == 0.0
    assert kw["width"] == float(prs.slide_width)
    assert kw["height"] == float(prs.slide_height)


def test_read_pptx_skips_evidence_when_dep_missing(monkeypatch):
    """The existing no-python-pptx no-op degradation is unaffected by the wiring."""
    monkeypatch.setitem(sys.modules, "pptx", None)
    store = _FakeStore()
    monkeypatch.setattr(native_ingest, "media_store", lambda: store)
    out = ro.read_pptx("/data/missing-dep.pptx")
    assert out == ""
    assert store.page_calls == []
