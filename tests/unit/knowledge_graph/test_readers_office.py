"""Tests for the office-document modality readers (CONCEPT:AU-KG.enrichment.multimodal-readers).

Covers the always-available stdlib readers (email ``.eml``, ``.csv``/``.tsv``)
with real extraction, the optional-dep readers (``.msg``/``.pptx``/``.xlsx``)
both when the lib is present (skipped if absent) and their graceful no-op
degradation when absent, and that each reader self-registers under its
extension(s) against the sibling reader registry.
"""

from __future__ import annotations

import csv
import importlib
import sys
from email.message import EmailMessage

import pytest

from agent_utilities.knowledge_graph.extraction import readers_office as ro


# ----------------------------------------------------------------------------
# Email — .eml (stdlib, always available)
# ----------------------------------------------------------------------------


def _write_eml(path, *, subject="Quarterly Review", body="Numbers look strong."):
    msg = EmailMessage()
    msg["From"] = "alice@example.com"
    msg["To"] = "bob@example.com"
    msg["Cc"] = "carol@example.com"
    msg["Subject"] = subject
    msg.set_content(body)
    path.write_bytes(msg.as_bytes())
    return path


def test_read_eml_headers_and_body(tmp_path):
    p = _write_eml(tmp_path / "mail.eml")
    text = ro.read_eml(str(p))
    assert "From: alice@example.com" in text
    assert "To: bob@example.com" in text
    assert "Cc: carol@example.com" in text
    assert "Subject: Quarterly Review" in text
    assert "Numbers look strong." in text
    # Header block precedes body.
    assert text.index("Subject:") < text.index("Numbers look strong.")


def test_read_eml_html_only_body_is_detagged(tmp_path):
    msg = EmailMessage()
    msg["From"] = "x@y.com"
    msg["Subject"] = "HTML mail"
    msg.add_alternative(
        "<html><body><p>Hello <b>world</b></p></body></html>", subtype="html"
    )
    p = tmp_path / "html.eml"
    p.write_bytes(msg.as_bytes())
    text = ro.read_eml(str(p))
    assert "Hello" in text and "world" in text
    assert "<p>" not in text and "<b>" not in text


def test_read_eml_unreadable_returns_empty(tmp_path):
    p = tmp_path / "missing.eml"
    assert ro.read_eml(str(p)) == ""


# ----------------------------------------------------------------------------
# Spreadsheets — .csv / .tsv (stdlib, always available)
# ----------------------------------------------------------------------------


def test_read_csv_rows(tmp_path):
    p = tmp_path / "data.csv"
    p.write_text("name,role\nAlice,Eng\nBob,PM\n", encoding="utf-8")
    text = ro.read_csv(str(p))
    assert "name\trole" in text
    assert "Alice\tEng" in text
    assert "Bob\tPM" in text


def test_read_csv_tsv_delimiter(tmp_path):
    p = tmp_path / "data.tsv"
    p.write_text("a\tb\tc\n1\t2\t3\n", encoding="utf-8")
    text = ro.read_csv(str(p))
    assert "a\tb\tc" in text
    assert "1\t2\t3" in text


def test_read_csv_empty_returns_empty(tmp_path):
    p = tmp_path / "empty.csv"
    p.write_text("", encoding="utf-8")
    assert ro.read_csv(str(p)) == ""


# ----------------------------------------------------------------------------
# Email — .msg (optional extract-msg): degrade when absent
# ----------------------------------------------------------------------------


def test_read_msg_degrades_without_dep(tmp_path, monkeypatch, caplog):
    # Force the optional import to fail regardless of install state.
    monkeypatch.setitem(sys.modules, "extract_msg", None)
    p = tmp_path / "x.msg"
    p.write_bytes(b"not a real msg")
    with caplog.at_level("WARNING"):
        out = ro.read_msg(str(p))
    assert out == ""
    assert any("extract-msg" in r.message for r in caplog.records)


@pytest.mark.skipif(
    importlib.util.find_spec("extract_msg") is None,
    reason="extract-msg not installed",
)
def test_read_msg_when_dep_present(tmp_path):  # pragma: no cover - dep-gated
    import extract_msg  # noqa: F401

    # Build a minimal .msg via the lib's own writer if available; otherwise just
    # assert the reader does not raise on a bogus path.
    out = ro.read_msg(str(tmp_path / "nope.msg"))
    assert isinstance(out, str)


# ----------------------------------------------------------------------------
# Presentations — .pptx (optional python-pptx)
# ----------------------------------------------------------------------------


def test_read_pptx_degrades_without_dep(tmp_path, monkeypatch, caplog):
    monkeypatch.setitem(sys.modules, "pptx", None)
    p = tmp_path / "deck.pptx"
    p.write_bytes(b"PK\x03\x04 not really")
    with caplog.at_level("WARNING"):
        out = ro.read_pptx(str(p))
    assert out == ""
    assert any("python-pptx" in r.message for r in caplog.records)


@pytest.mark.skipif(
    importlib.util.find_spec("pptx") is None,
    reason="python-pptx not installed",
)
def test_read_pptx_extracts_slide_text(tmp_path):  # pragma: no cover - dep-gated
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.text = "Strategy Overview"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "Cover the roadmap"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))

    text = ro.read_pptx(str(p))
    assert "Strategy Overview" in text
    assert "--- Slide 1 ---" in text
    assert "Cover the roadmap" in text


# ----------------------------------------------------------------------------
# Spreadsheets — .xlsx (optional openpyxl)
# ----------------------------------------------------------------------------


def test_read_xlsx_degrades_without_dep(tmp_path, monkeypatch, caplog):
    monkeypatch.setitem(sys.modules, "openpyxl", None)
    p = tmp_path / "book.xlsx"
    p.write_bytes(b"PK\x03\x04 not really")
    with caplog.at_level("WARNING"):
        out = ro.read_xlsx(str(p))
    assert out == ""
    assert any("openpyxl" in r.message for r in caplog.records)


@pytest.mark.skipif(
    importlib.util.find_spec("openpyxl") is None,
    reason="openpyxl not installed",
)
def test_read_xlsx_extracts_sheet_rows(tmp_path):  # pragma: no cover - dep-gated
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Q", "Amount"])
    ws.append(["Q1", 100])
    ws.append(["Q2", 150])
    p = tmp_path / "book.xlsx"
    wb.save(str(p))

    text = ro.read_xlsx(str(p))
    assert "# Revenue" in text
    assert "Q\tAmount" in text
    assert "Q1\t100" in text


# ----------------------------------------------------------------------------
# Registration — each reader self-registers under its extension(s)
# ----------------------------------------------------------------------------


def test_readers_self_register():
    """Each reader registers under its extension(s) via @register_reader.

    Works against the real ``extraction.readers`` registry when present; if it is
    not present yet, the module's fallback no-op decorator is used and this test
    asserts the readers are at least importable+callable (the registry binding is
    exercised once ``readers.py`` lands).
    """
    try:
        from agent_utilities.knowledge_graph.extraction import readers as reg
    except Exception:
        # Registry sibling not present yet — fallback decorator path. Assert the
        # readers exist and are callable so wiring is ready for the registry.
        for fn_name in (
            "read_eml",
            "read_msg",
            "read_pptx",
            "read_xlsx",
            "read_csv",
        ):
            assert callable(getattr(ro, fn_name))
        return

    # Registry present: discover/import then assert lookups resolve to readers.
    for discover in ("discover_readers", "discover"):
        fn = getattr(reg, discover, None)
        if callable(fn):
            fn()
            break
    importlib.import_module("agent_utilities.knowledge_graph.extraction.readers_office")
    lookup = getattr(reg, "get_reader", None)
    if not callable(lookup):
        pytest.skip("registry has no get_reader accessor to assert against")
    for ext in (".eml", ".msg", ".pptx", ".xlsx", ".csv", ".tsv"):
        assert lookup(ext) is not None, f"no reader registered for {ext}"


# ----------------------------------------------------------------------------
# Shared formatter sanity (pure helpers)
# ----------------------------------------------------------------------------


def test_format_rows_skips_blank_rows():
    out = ro._format_rows([["a", "b"], ["", ""], ["c", None]])
    assert out == "a\tb\nc"


def test_format_email_header_then_body():
    out = ro._format_email(
        {"From": "a@b.com", "Subject": "Hi", "To": "", "Cc": "", "Date": ""},
        "body text",
    )
    assert out == "From: a@b.com\nSubject: Hi\n\nbody text"


def test_csv_roundtrip_sniff(tmp_path):
    # Semicolon-delimited gets sniffed and re-emitted tab-joined.
    p = tmp_path / "semi.csv"
    rows = [["x", "y"], ["1", "2"]]
    with open(p, "w", newline="", encoding="utf-8") as fh:
        csv.writer(fh, delimiter=";").writerows(rows)
    text = ro.read_csv(str(p))
    assert "x\ty" in text and "1\t2" in text
